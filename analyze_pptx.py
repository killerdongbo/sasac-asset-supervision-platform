"""Analyze PPTX structure to understand slide layouts and design."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
import json

pptx_path = r"d:\谢东波个人\智算\国资监管平台\国有资产监管系统升级项目汇报-V0.4.pptx"
prs = Presentation(pptx_path)

print(f"Slide width: {prs.slide_width}, height: {prs.slide_height}")
print(f"Slide width (inches): {prs.slide_width/914400}, height (inches): {prs.slide_height/914400}")
print(f"Number of slides: {len(prs.slides)}")
print(f"Number of slide layouts: {len(prs.slide_layouts)}")
print(f"Number of slide masters: {len(prs.slide_masters)}")

print("\n=== SLIDE LAYOUTS ===")
for i, layout in enumerate(prs.slide_layouts):
    print(f"\nLayout {i}: {layout.name}")
    print(f"  Placeholders: {len(layout.placeholders)}")
    for ph in layout.placeholders:
        print(f"    - {ph.placeholder_format.idx}: {ph.name} (type={ph.placeholder_format.type}, {ph.width}x{ph.height})")

print("\n=== CHAPTER 3 SLIDES (17-25) ===")
for i in range(16, 25):  # slides 17-25 (0-indexed: 16-24)
    if i >= len(prs.slides):
        break
    slide = prs.slides[i]
    print(f"\n--- Slide {i+1} ---")
    print(f"  Layout: {slide.slide_layout.name}")
    shapes_info = []
    for j, shape in enumerate(slide.shapes):
        info = {
            'id': j,
            'name': shape.name,
            'type': str(shape.shape_type),
            'left': shape.left,
            'top': shape.top,
            'width': shape.width,
            'height': shape.height,
        }
        if shape.has_text_frame:
            info['text'] = shape.text_frame.text[:100] if shape.text_frame.text else ''
        shapes_info.append(info)
        text_preview = ''
        if shape.has_text_frame:
            text_preview = shape.text_frame.text[:80].replace('\n', '|')
        print(f"  Shape {j}: {shape.name} ({shape.shape_type}) @ ({shape.left}, {shape.top}) {shape.width}x{shape.height} | '{text_preview}'")

print("\n=== SLIDE MASTERS ===")
for i, master in enumerate(prs.slide_masters):
    print(f"\nMaster {i}: {master.name}")
    print(f"  Slide layouts: {len(master.slide_layouts)}")
    for j, layout in enumerate(master.slide_layouts):
        print(f"  Layout {j}: {layout.name}")
