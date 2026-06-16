"""
Add highlight slides to Chapter 3 of the PPTX.
Strategy: Clone an existing content slide (slide 19 has a clean grid layout) via XML manipulation,
then modify text content to create new highlight slides inserted before slide 25 (Chapter 4).
"""
import sys
import io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8')

import os
import copy
import shutil
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.oxml.ns import qn, nsmap
import re

# Paths
SRC_PATH = r"d:\谢东波个人\智算\国资监管平台\国有资产监管系统升级项目汇报-V0.4.pptx"
DST_PATH = r"d:\谢东波个人\智算\国资监管平台\国有资产监管系统升级项目汇报-V0.5.pptx"

# Copy the original file
print("Copying original file...")
shutil.copy2(SRC_PATH, DST_PATH)

# Open the copy
prs = Presentation(DST_PATH)

# Register all namespaces
nsmap2 = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'p14': 'http://schemas.microsoft.com/office/powerpoint/2010/main',
}

def clone_slide(prs, source_slide_idx):
    """
    Clone a slide and return the new slide's XML element.
    Uses the standard python-pptx workaround for slide cloning.
    """
    # Get source slide data
    source_slide = prs.slides[source_slide_idx]

    # Get the source slide's relationships
    source_slide_part = source_slide.part
    source_rel = source_slide_part.rels

    # Get source layout
    source_layout = source_slide.slide_layout

    # Add a new slide with the same layout
    # python-pptx adds slides to the END, we'll reorder later
    blank_slide_layout = prs.slide_layouts[7]  # 'Blank' layout
    new_slide = prs.slides.add_slide(blank_slide_layout)

    # Get new slide XML
    new_slide_part = new_slide.part

    # Copy the source slide XML to the new slide
    # First remove existing shapes from new slide
    spTree = new_slide.shapes._spTree
    for child in list(spTree):
        if child.tag != qn('p:nvGrpSpPr'):
            spTree.remove(child)

    # Copy shapes from source to new slide
    source_spTree = source_slide.shapes._spTree
    for child in list(source_spTree):
        if child.tag != qn('p:nvGrpSpPr'):
            imported = copy.deepcopy(child)
            spTree.append(imported)

    # Copy image relationships from source
    imported_rids = {}
    for rel in source_rel.values():
        if "image" in rel.reltype or "media" in rel.reltype:
            # Handle relationship copying for images
            try:
                new_rid = new_slide_part.relate_to(rel.target_part, rel.reltype)
                imported_rids[rel.rId] = new_rid
            except Exception as e:
                print(f"  Warning: Could not copy relationship {rel.rId}: {e}")

    # Update rId references in the new slide
    if imported_rids:
        for elem in spTree.iter():
            for attr_name in ['{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed',
                            '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}link']:
                old_rid = elem.get(attr_name)
                if old_rid and old_rid in imported_rids:
                    elem.set(attr_name, imported_rids[old_rid])

    # Also need to update the slide layout reference to match source
    new_slide_part.slide_layout = source_slide_part.slide_layout

    return new_slide, new_slide_part

def clone_and_insert_slides(prs):
    """
    Clone slide 18 (index 17) which has the 3-box architecture design,
    and create 5 new highlight slides inserted before slide 25 (index 24, Chapter 4 title).
    """
    # The source template slide - slide 19 (index 18) has clean rounded rectangle layout
    # But slide 18 (index 17) has the more impressive 3-box design
    # Let's use slide 18 as the base template

    template_idx = 17  # Slide 18 (0-indexed) - 系统总体架构

    highlights = [
        {
            'title': '资产粒子化管理：一物一码·权属重划',
            'subtitle': '每项资产赋予唯一数字编码，实现资产权属的灵活重新划分与精准追溯',
            'blocks': [
                ('01', '资产唯一编码', '每项资产生成全局唯一编码（AssetCode），支持QR码/条码双模式输出，实现\"一物一码\"精准定位'),
                ('02', '权属灵活重划', '支持资产归属权的在线调整与重新划分，适配企业重组、资产划转等国资改革场景，权属变更全程留痕'),
                ('03', '全链精准追溯', '从采购入账到报废退出，每项资产全生命周期均可通过唯一编码一键追溯，操作记录不可篡改'),
                ('03B', 'AI辅助估值', '结合历史折旧数据与市场行情，AI模型辅助资产重估定价，为权属划转提供科学依据'),
            ],
            'footer': '核心亮点：一物一码打通资产\"身份证\"体系  |  权属变更流程化、可审计  |  为国资穿透式监管提供数据基础'
        },
        {
            'title': '数据处理架构：多源汇聚·流批一体',
            'subtitle': '构建企业级数据中台，支撑国资海量资产的实时采集、清洗、计算与分析',
            'blocks': [
                ('01', '多源数据接入', '适配企业现有ERP/EAM/财务系统，支持数据库直连、API对接、文件批量导入等多种数据源接入方式'),
                ('02', '流批一体引擎', '基于Flink/Spark构建流批一体数据处理引擎，实时同步资产变更事件，批量计算资产折旧与报表'),
                ('03B', '多层级数据治理', '数据质量校验规则引擎：完整性/一致性/准确性自动校验；异常数据自动识别并生成修正工单'),
                ('03', '国产化适配', '全面适配达梦/人大金仓等国产数据库，支持ARM/x86混合架构部署，信创环境下数据处理性能无损'),
            ],
            'footer': '技术选型：SpringBoot3 + Flink/Spark + PostgreSQL + Redis  |  数据治理：DQC规则引擎自动校验  |  国产数据库全面兼容'
        },
        {
            'title': '资产穿透式统计：逐级汇总·层层穿透',
            'subtitle': '支持从国资委→集团→企业→部门→单资产的多层级穿透式数据统计与钻取分析',
            'blocks': [
                ('01', '五级穿透查询', '国资委→集团→企业→部门→单项资产，五级穿透式查询，逐层展开资产总值、分类构成与变动趋势'),
                ('02', '分层统计引擎', '按组织层级预计算资产总值、净值、折旧额等核心指标，支持按分类/状态/时间等多维度交叉分析'),
                ('03', '实时驾驶舱', '管理驾驶舱可视呈现各层级资产总值分布、同比环比变化、闲置率等关键指标，一屏掌控全局'),
            ],
            'footer': '关键能力：5级组织穿透  |  KPI预计算加速  |  多维度交叉分析  |  一屏掌控全局资产格局'
        },
        {
            'title': '智能预警与全链路审计追溯',
            'subtitle': '主动式风险预警 + 不可篡改的审计追溯，构建国资监管安全防线',
            'blocks': [
                ('01', '多维度智能预警', '维保到期/巡检超期/资产闲置/折旧异常等多维度预警规则，定时自动扫描，主动推送通知'),
                ('02', '全链路审计追溯', '基于JSONB快照技术记录每次操作的before/after数据，审计日志独立存储且不可篡改，满足国资委审计要求'),
                ('03', '风险雷达图', '对集团下各企业资产健康度进行综合评分，生成风险雷达图，辅助监管决策与资源调配'),
            ],
            'footer': '安全底线：事前预警→事中阻断→事后可溯  |  JSONB数据快照+独立审计库  |  风险可视化辅助决策'
        },
        {
            'title': '数据资产入表与价值释放',
            'subtitle': '响应国家\"数据二十条\"，为国资数据资产入表提供全流程支撑',
            'blocks': [
                ('01', '数据资产盘点', '自动识别可入表的数据资产类型，建立数据资产目录，评估数据资产质量等级'),
                ('02', '入表流程支撑', '数据资产确认→成本归集→价值评估→会计入表→后续计量，全流程线上化管理'),
                ('03', '价值释放通道', '安全脱敏后的资产统计数据支持对外授权运营，探索国资数据要素市场化配置路径'),
            ],
            'footer': '政策响应：\"数据二十条\"落地  |  数据资产全流程管理  |  国资数据要素价值释放'
        },
    ]

    print(f"\nCloning template slide (index {template_idx}: '{prs.slides[template_idx].shapes[0].text_frame.text[:30] if prs.slides[template_idx].shapes[0].has_text_frame else 'N/A'}')...")
    print(f"Creating {len(highlights)} highlight slides...")

    created_slides = []
    for i, hl in enumerate(highlights):
        print(f"\n--- Creating highlight slide {i+1}: {hl['title'][:40]}... ---")
        new_slide, new_slide_part = clone_slide(prs, template_idx)

        # Now modify the text content
        shapes = new_slide.shapes
        for shape in shapes:
            if shape.has_text_frame:
                tf = shape.text_frame
                full_text = tf.text

                # Title text box
                if '系统总体架构' in full_text:
                    tf.clear()
                    p = tf.paragraphs[0]
                    run = p.add_run()
                    run.text = hl['title']
                    # Copy font properties from original
                    for para in shape.text_frame.paragraphs:
                        for r in para.runs:
                            if r.font.size:
                                run.font.size = r.font.size
                            break

                # Subtitle/description text boxes
                elif '三域模块化设计' in full_text or ('三域' in full_text and '设计' in full_text):
                    # This is a decorative group - keep it
                    pass

                # The three domain blocks - modify them
                elif '企业资产管理域' in full_text or '国资监管域' in full_text or '平台基础域' in full_text:
                    # We'll handle these via shape name matching instead
                    pass

        # More targeted text replacement using shape iteration
        shape_list = list(shapes)
        for shape in shape_list:
            if shape.has_text_frame:
                tf = shape.text_frame
                text = tf.text

                # Update the title text box (contains "系统总体架构")
                if '系统总体架构：三域模块化设计' in text:
                    tf.clear()
                    p = tf.paragraphs[0]
                    p.alignment = 1  # center
                    run = p.add_run()
                    run.text = hl['title']
                    # Try to match font size (36pt = 2592000 EMU, or 18pt for space constraints)
                    run.font.size = Pt(32)
                    run.font.bold = True
                    run.font.color.rgb = None  # will inherit from shape
                    try:
                        rPr = run._r.get_or_add_rPr()
                        solidFill = etree.SubElement(rPr, qn('a:solidFill'))
                        srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
                        srgbClr.set('val', 'FFFFFF')
                    except:
                        pass
                    print(f"  Updated title: '{hl['title']}'")
                    continue

                # Update the subtitle / description text
                if '三域模块化设计' in text:
                    tf.clear()
                    p = tf.paragraphs[0]
                    run = p.add_run()
                    run.text = hl['subtitle']
                    run.font.size = Pt(18)
                    try:
                        rPr = run._r.get_or_add_rPr()
                        solidFill = etree.SubElement(rPr, qn('a:solidFill'))
                        srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
                        srgbClr.set('val', 'FFFFFF')
                    except:
                        pass
                    continue

                # Update block content - match by block number patterns
                # The content blocks have specific patterns
                if '核心：资产全生命周期管理' in text:
                    self._set_block_text(tf, hl['blocks'][0][1], hl['blocks'][0][2])
                    continue
                if '核心：监管决策与风险防控' in text:
                    self._set_block_text(tf, hl['blocks'][1][1], hl['blocks'][1][2])
                    continue
                if '横切能力：统一提供，自动继承' in text:
                    self._set_block_text(tf, hl['blocks'][2][1], hl['blocks'][2][2])
                    continue

                # Update footer
                if '关键设计决策' in text:
                    self._set_footer_text(tf, hl['footer'])
                    continue

        # Also handle domain name labels
        for shape in shape_list:
            if shape.has_text_frame:
                tf = shape.text_frame
                text = tf.text
                if text == '企业资产管理域':
                    self._set_simple_text(tf, hl['blocks'][0][0] if len(hl['blocks']) > 3 else '')
                elif text == '国资监管域':
                    self._set_simple_text(tf, hl['blocks'][1][0] if len(hl['blocks']) > 3 else '')
                elif text == '平台基础域':
                    self._set_simple_text(tf, hl['blocks'][2][0] if len(hl['blocks']) > 3 else '')

        created_slides.append((new_slide, new_slide_part))
        print(f"  Done creating slide: {hl['title'][:40]}")

    return created_slides


def _set_block_text(tf, title, content):
    """Set block text with title and content."""
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title + '\n\n' + content
    run.font.size = Pt(12)
    try:
        rPr = run._r.get_or_add_rPr()
        solidFill = etree.SubElement(rPr, qn('a:solidFill'))
        srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
        srgbClr.set('val', 'FFFFFF')
    except:
        pass


def _set_footer_text(tf, text):
    """Set footer text."""
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = 1
    run = p.add_run()
    run.text = text
    run.font.size = Pt(11)
    try:
        rPr = run._r.get_or_add_rPr()
        solidFill = etree.SubElement(rPr, qn('a:solidFill'))
        srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
        srgbClr.set('val', '000000')
    except:
        pass


def _set_simple_text(tf, text):
    """Set simple label text."""
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = 1
    run = p.add_run()
    run.text = text
    run.font.size = Pt(18)
    try:
        rPr = run._r.get_or_add_rPr()
        solidFill = etree.SubElement(rPr, qn('a:solidFill'))
        srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
        srgbClr.set('val', 'FFFFFF')
    except:
        pass


def reorder_slides(prs, created_slides):
    """
    Reorder slides: move newly created slides to appear after slide 24 (index 23)
    and before slide 25 (Chapter 4 title slide).

    Current layout after adding slides (at end):
    [0..23] existing slides (slide 1-24), [24] old slide 25 (Ch4), [25..34] old slides 26-33, [35..39] new slides

    Target:
    [0..23] existing, [24..28] new highlight slides, [29] old slide 25 (Ch4), [30..] old slides 26-33
    """
    # Get the presentation XML
    pres_part = prs.part
    pres_xml = pres_part.element

    # Find the sldIdLst element
    sldIdLst = pres_xml.find(qn('p:sldIdLst'))
    if sldIdLst is None:
        print("ERROR: Could not find sldIdLst in presentation.xml")
        return

    # Get all sldId elements
    sld_ids = list(sldIdLst.findall(qn('p:sldId')))
    print(f"\nTotal slides in presentation: {len(sld_ids)}")

    # The new slides were added at the end (last 5)
    # Slide index 24 (0-indexed) is the Chapter 4 title slide
    # We need to move the last 5 slides to after index 23

    num_new = len(created_slides)
    total = len(sld_ids)

    # The new slides are the last N entries
    new_sld_ids = sld_ids[-num_new:]
    old_sld_ids = sld_ids[:-num_new]

    # Insert new slides after position 23 (after slide 24, before slide 25 which is Chapter 4)
    insert_pos = 24  # after slide 24 (0-indexed), before the Chapter 4 title slide

    # Reorder: [0..23] + new_slides + [24..end]
    reordered = old_sld_ids[:insert_pos] + new_sld_ids + old_sld_ids[insert_pos:]

    # Clear and re-add in new order
    for sld_id in list(sldIdLst):
        sldIdLst.remove(sld_id)

    for sld_id in reordered:
        sldIdLst.append(sld_id)

    print(f"Reordered slides: {len(reordered)} total")
    print(f"New highlight slides inserted at positions {insert_pos+1}-{insert_pos+num_new} (1-indexed: slides {insert_pos+2}-{insert_pos+num_new+1})")


# Main execution
print("=" * 60)
print("Adding highlight slides to Chapter 3...")
print("=" * 60)

# Clone and modify slides
created = clone_and_insert_slides(prs)

# Reorder to place them in Chapter 3
reorder_slides(prs, created)

# Save the modified presentation
print("\nSaving modified presentation...")
prs.save(DST_PATH)
print(f"\nDone! Output saved to: {DST_PATH}")
print(f"Original: {len(prs.slides)} slides")
print(f"Added {len(created)} new highlight slides in Chapter 3")
