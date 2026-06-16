"""Analyze PPTX structure properly with UTF-8."""
import sys
import io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor

pptx_path = r"d:\谢东波个人\智算\国资监管平台\国有资产监管系统升级项目汇报-V0.4.pptx"
prs = Presentation(pptx_path)

print(f"Slides: {len(prs.slides)}, Width: {prs.slide_width}, Height: {prs.slide_height}")

print("\n=== SLIDE LAYOUTS ===")
for i, layout in enumerate(prs.slide_layouts):
    ph_names = [ph.name for ph in layout.placeholders]
    print(f"Layout {i}: '{layout.name}' placeholders={len(layout.placeholders)} names={ph_names}")

print("\n=== CHAPTER 3 SLIDES (17-25) ===")
for slide_idx in range(16, 25):
    if slide_idx >= len(prs.slides):
        break
    slide = prs.slides[slide_idx]
    print(f"\n--- Slide {slide_idx+1} (layout='{slide.slide_layout.name}') ---")
    for j, shape in enumerate(slide.shapes):
        text = ''
        if shape.has_text_frame:
            text = shape.text_frame.text.replace('\n', '\\n').replace('\r', '')[:120]
        print(f"  [{j}] {shape.shape_type} '{shape.name}' pos=({shape.left},{shape.top}) size=({shape.width},{shape.height}) text='{text}'")
