"""
Add highlight slides to Chapter 3 of the PPTX.
Uses python-pptx to create new slides with matching color scheme and layout style,
then reorders them within Chapter 3.
"""
import sys
import io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8')

import os
import copy
import shutil
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

# ========== CONSTANTS ==========
SRC_PATH = r"d:\谢东波个人\智算\国资监管平台\国有资产监管系统升级项目汇报-V0.4.pptx"
DST_PATH = r"d:\谢东波个人\智算\国资监管平台\国有资产监管系统升级项目汇报-V0.5.pptx"

# Color scheme matching the existing design
ACCENT_BLUE = RGBColor(0x3C, 0x89, 0xFF)  # Primary blue
DARK_BLUE = RGBColor(0x1A, 0x56, 0xD6)    # Darker blue for contrast
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x1A)
MEDIUM_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_BLUE_BG = RGBColor(0xE8, 0xF0, 0xFF)
GRADIENT_GREEN = RGBColor(0x2A, 0xF5, 0x98)
GRADIENT_CYAN = RGBColor(0x08, 0xAE, 0xEA)
ORANGE_ACCENT = RGBColor(0xFF, 0x8C, 0x00)

# Slide dimensions (standard 16:9)
SLIDE_W = 12192000  # EMU
SLIDE_H = 6858000   # EMU

# Font settings (match existing: 汉仪润圆-65简 as primary, 微软雅黑 fallback)
TITLE_FONT = '微软雅黑'
BODY_FONT = '微软雅黑'

def hex_to_rgb(hex_str):
    """Convert hex color string to RGBColor."""
    hex_str = hex_str.lstrip('#')
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))


def add_title_bar(slide, title_text, subtitle_text=''):
    """Add a top title area matching the existing slide style: title + decorative lines."""
    # Main title
    txBox = slide.shapes.add_textbox(Cm(8), Cm(0.5), Cm(18), Cm(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.name = TITLE_FONT
    run.font.color.rgb = DARK_BLUE

    # Subtitle if provided
    if subtitle_text:
        txBox2 = slide.shapes.add_textbox(Cm(6), Cm(1.6), Cm(22), Cm(0.6))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = subtitle_text
        run2.font.size = Pt(14)
        run2.font.name = BODY_FONT
        run2.font.color.rgb = MEDIUM_GRAY

    # Decorative line under title
    line = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Cm(12), Cm(1.35), Cm(10), Cm(0.03)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()

    # Short decoration lines on sides
    for x_cm in [9, 23]:
        side_line = slide.shapes.add_shape(
            1, Cm(x_cm), Cm(1.35), Cm(3), Cm(0.02)
        )
        side_line.fill.solid()
        side_line.fill.fore_color.rgb = ACCENT_BLUE
        side_line.line.fill.background()


def add_footer_bar(slide, footer_text):
    """Add a footer bar at the bottom matching the existing style."""
    # Background bar
    bar = slide.shapes.add_shape(
        1, Cm(0.5), Cm(17.5), Cm(32.5), Cm(0.8)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0xF0, 0xF5, 0xFF)
    bar.line.fill.background()

    # Footer text
    txBox = slide.shapes.add_textbox(Cm(1), Cm(17.55), Cm(31.5), Cm(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = footer_text
    run.font.size = Pt(10)
    run.font.name = BODY_FONT
    run.font.color.rgb = MEDIUM_GRAY


def add_feature_card(slide, left_cm, top_cm, width_cm, height_cm, icon_text, title, description, color=None):
    """Add a rounded rectangle card with icon, title and description - matching the existing card style."""
    if color is None:
        color = ACCENT_BLUE

    # Card background (rounded rectangle)
    card = slide.shapes.add_shape(
        5,  # MSO_SHAPE.ROUNDED_RECTANGLE
        Cm(left_cm), Cm(top_cm), Cm(width_cm), Cm(height_cm)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = color
    card.line.fill.background()
    # Set corner radius
    card.adjustments[0] = 0.1

    # Icon circle
    icon_size = min(1.2, width_cm * 0.25)
    icon = slide.shapes.add_shape(
        9,  # MSO_SHAPE.OVAL
        Cm(left_cm + 0.5), Cm(top_cm + 0.4), Cm(icon_size), Cm(icon_size)
    )
    icon.fill.solid()
    icon.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    icon.line.fill.background()

    # Icon text (number)
    tf_icon = icon.text_frame
    tf_icon.word_wrap = False
    p_icon = tf_icon.paragraphs[0]
    p_icon.alignment = PP_ALIGN.CENTER
    run_icon = p_icon.add_run()
    run_icon.text = icon_text[:2]
    run_icon.font.size = Pt(14)
    run_icon.font.bold = True
    run_icon.font.name = BODY_FONT
    run_icon.font.color.rgb = color

    # Title text
    txBox_title = slide.shapes.add_textbox(
        Cm(left_cm + 0.5), Cm(top_cm + 0.3), Cm(width_cm - 1.0), Cm(0.8)
    )
    tf_title = txBox_title.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.alignment = PP_ALIGN.LEFT
    run_title = p_title.add_run()
    run_title.text = icon_text + '  ' + title
    run_title.font.size = Pt(13)
    run_title.font.bold = True
    run_title.font.name = BODY_FONT
    run_title.font.color.rgb = WHITE

    # Description text
    desc_top = top_cm + 1.5
    txBox_desc = slide.shapes.add_textbox(
        Cm(left_cm + 0.5), Cm(desc_top), Cm(width_cm - 1.0), Cm(height_cm - desc_top + top_cm - 0.4)
    )
    tf_desc = txBox_desc.text_frame
    tf_desc.word_wrap = True
    p_desc = tf_desc.paragraphs[0]
    p_desc.alignment = PP_ALIGN.LEFT
    p_desc.line_spacing = Pt(18)
    run_desc = p_desc.add_run()
    run_desc.text = description
    run_desc.font.size = Pt(10)
    run_desc.font.name = BODY_FONT
    run_desc.font.color.rgb = WHITE


def create_highlight_slide(prs, highlight_data):
    """Create a single highlight slide with cards layout."""
    # Use Layout 4 (仅标题) which has no placeholders - clean content slide
    layout = prs.slide_layouts[4]  # '仅标题'
    slide = prs.slides.add_slide(layout)

    # White background
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE

    # Add title bar
    add_title_bar(slide, highlight_data['title'], highlight_data.get('subtitle', ''))

    # Add feature cards
    blocks = highlight_data['blocks']
    num_blocks = len(blocks)

    if num_blocks <= 3:
        # 3 cards in a row
        card_w = 10.0
        spacing = 0.8
        total_w = num_blocks * card_w + (num_blocks - 1) * spacing
        start_x = (33.8 - total_w) / 2
        card_h = 9.0

        for i, (icon, title, desc) in enumerate(blocks):
            x = start_x + i * (card_w + spacing)
            colors = [ACCENT_BLUE, RGBColor(0x00, 0x96, 0x88), RGBColor(0xFF, 0x6F, 0x00), RGBColor(0x6A, 0x1B, 0x9A)]
            add_feature_card(slide, x, 3.2, card_w, card_h, icon, title, desc, colors[i % len(colors)])

    elif num_blocks == 4:
        # 2x2 grid
        card_w = 15.0
        card_h = 5.5
        spacing_x = 1.0
        spacing_y = 0.6
        start_x = (33.8 - 2 * card_w - spacing_x) / 2
        start_y = 3.0

        for i, (icon, title, desc) in enumerate(blocks):
            col = i % 2
            row = i // 2
            x = start_x + col * (card_w + spacing_x)
            y = start_y + row * (card_h + spacing_y)
            colors = [ACCENT_BLUE, RGBColor(0x00, 0x96, 0x88), RGBColor(0xFF, 0x6F, 0x00), RGBColor(0x6A, 0x1B, 0x9A)]
            add_feature_card(slide, x, y, card_w, card_h, icon, title, desc, colors[i % len(colors)])

    elif num_blocks == 5:
        # Top row: 3 cards, bottom row: 2 cards
        card_w = 10.0
        card_h = 5.0
        spacing_x = 0.8
        spacing_y = 0.6

        # Top row: 3 cards
        for i in range(3):
            total_w_top = 3 * card_w + 2 * spacing_x
            start_x = (33.8 - total_w_top) / 2
            x = start_x + i * (card_w + spacing_x)
            y = 3.0
            colors = [ACCENT_BLUE, RGBColor(0x00, 0x96, 0x88), RGBColor(0xFF, 0x6F, 0x00), RGBColor(0x6A, 0x1B, 0x9A)]
            icon, title, desc = blocks[i]
            add_feature_card(slide, x, y, card_w, card_h, icon, title, desc, colors[i % len(colors)])

        # Bottom row: 2 cards, centered
        for i in range(2):
            total_w_bottom = 2 * card_w + spacing_x
            start_x = (33.8 - total_w_bottom) / 2
            x = start_x + i * (card_w + spacing_x)
            y = 3.0 + card_h + spacing_y
            icon, title, desc = blocks[3 + i]
            colors = [RGBColor(0x6A, 0x1B, 0x9A), RGBColor(0xE9, 0x1E, 0x63)]
            add_feature_card(slide, x, y, card_w, card_h, icon, title, desc, colors[i % len(colors)])

    # Add footer
    if 'footer' in highlight_data:
        add_footer_bar(slide, highlight_data['footer'])

    return slide


# ========== HIGHLIGHT CONTENT ==========
HIGHLIGHTS = [
    {
        'title': '核心亮点一：资产粒子化管理 —— 一物一码·权属重划',
        'subtitle': '每项资产赋予唯一数字编码，支持资产归属权的灵活重新划分与全生命周期精准追溯',
        'blocks': [
            ('01', '资产唯一编码体系',
             '每项资产自动生成全局唯一编码，支持QR码与条码双模式输出，实现"一物一码"精准定位，扫码即可查看资产全生命周期信息。'),
            ('02', '资产权属灵活重划',
             '支持资产归属组织在线调整与重新划分，适配国企重组、资产划转等改革场景，权属变更全流程留痕、可审计、可追溯。'),
            ('03', '全生命周期精准追溯',
             '从采购入账到报废退出，每项资产全生命周期均通过唯一编码一键追溯，操作记录基于JSONB快照独立存储、不可篡改。'),
            ('04', 'AI辅助资产估值',
             '结合历史折旧数据、同类资产市场行情，AI模型辅助资产重估定价，为权属划转、资产处置提供科学价值依据。'),
        ],
        'footer': '核心能力：一物一码打通资产"身份证"体系  |  权属变更流程化、可审计、可回溯  |  为国资穿透式监管奠定数据基础'
    },
    {
        'title': '核心亮点二：企业级数据处理架构',
        'subtitle': '构建多源汇聚、流批一体的数据处理中台，支撑海量国资资产的实时采集、清洗、计算与分析',
        'blocks': [
            ('01', '多源数据汇聚接入',
             '适配企业现有ERP/EAM/财务系统，支持数据库直连、REST API对接、Excel批量导入等多种数据源接入方式，实现资产数据一站式汇聚。'),
            ('02', '流批一体计算引擎',
             '基于流批一体架构实现资产变更事件实时同步与批量折旧计算双模处理——变更数据秒级响应，报表计算T+1自动产出。'),
            ('03', '多层级数据质量治理',
             '内置数据质量校验规则引擎：完整性/一致性/准确性自动校验；异常数据自动识别并生成修正工单，确保监管数据零差错。'),
            ('04', '全国产化信创适配',
             '全面适配达梦DM8/人大金仓/神通等国产数据库，支持ARM/x86混合架构，信创环境下数据处理性能零损耗。'),
        ],
        'footer': '技术选型：SpringBoot3 + Flink/Spark双引擎 + PostgreSQL + Redis + MinIO  |  数据治理：DQC规则引擎全自动校验  |  国产数据库全面兼容'
    },
    {
        'title': '核心亮点三：资产穿透式统计 —— 逐级汇总·层层钻取',
        'subtitle': '支持从国资委到单资产的五级穿透式数据统计，实时掌握各层级资产总值、结构与变动趋势',
        'blocks': [
            ('01', '五级穿透查询',
             '国资委→集团→企业→部门→单项资产，五级穿透式查询，逐层展开资产总值、净值、折旧额与分类构成，数据一目了然。'),
            ('02', 'KPI分层预计算引擎',
             '按组织层级预计算资产总值、净值、折旧额等核心KPI，基于PostgreSQL物化视图+Redis缓存实现毫秒级查询响应。'),
            ('03', '多维度交叉分析',
             '支持按资产分类、使用状态、地域分布、时间周期等多维度交叉钻取分析，自动生成同比环比趋势图表与异常波动标注。'),
            ('04', '管理驾驶舱可视化',
             '大屏驾驶舱直观呈现各级资产总值分布、闲置率排行、预警概览等关键指标，一屏掌控全局，辅助监管决策。'),
        ],
        'footer': '关键能力：5级组织穿透查询  |  KPI预计算毫秒级响应  |  多维度交叉分析+同比环比  |  大屏驾驶舱一屏掌控全局'
    },
    {
        'title': '核心亮点四：智能预警与全链路审计追溯',
        'subtitle': '主动式风险预警体系 + 不可篡改的全链路审计追溯，构建国资监管安全双防线',
        'blocks': [
            ('01', '多维度智能预警',
             '维保到期、巡检超期、资产闲置、折旧异常、价值异动等多维度预警规则，定时自动扫描，支持短信/邮件/系统消息多通道推送。'),
            ('02', '全链路审计追溯',
             '基于JSONB快照技术记录每次操作的before/after完整数据，审计日志独立存储、不可篡改，满足国资委审计检查要求。'),
            ('03', '资产健康度评估',
             '对集团下各企业资产健康度进行多维度综合评分，生成风险雷达图，直观呈现资产质量分布，辅助精准监管与资源调配。'),
        ],
        'footer': '安全防线：事前预警→事中阻断→事后可溯  |  JSONB数据快照+独立审计存储  |  资产健康度风险雷达图辅助决策'
    },
    {
        'title': '核心亮点五：数据资产入表与要素价值释放',
        'subtitle': '响应国家"数据二十条"政策，为国资数据资产确认、计量、入表提供全流程数字化支撑',
        'blocks': [
            ('01', '数据资产自动盘点',
             '自动识别可入表的数据资产类别（资产统计数据、经营分析数据、监管指标数据），建立分级分类的数据资产目录。'),
            ('02', '入表全流程管理',
             '数据资产确认→成本归集→价值评估→会计入表→后续计量，全流程线上化审批管理，符合《企业数据资源相关会计处理暂行规定》。'),
            ('03', '数据要素安全流通',
             '安全脱敏后的资产统计数据支持对外授权运营，建立数据使用审批与计量计费机制，探索国资数据要素市场化配置路径。'),
        ],
        'footer': '政策响应："数据二十条"落地实践  |  数据资产入表全流程管理  |  国资数据要素安全流通与价值释放'
    },
]


def reorder_slides(prs, num_new):
    """
    Reorder slides: move the last `num_new` slides to after slide 24 (Chapter 3 end),
    before slide 25 (Chapter 4 title).
    """
    pres_part = prs.part
    pres_xml = pres_part._element

    sldIdLst = pres_xml.find(qn('p:sldIdLst'))
    if sldIdLst is None:
        print("ERROR: Could not find sldIdLst!")
        return

    sld_ids = list(sldIdLst.findall(qn('p:sldId')))
    total = len(sld_ids)

    # New slides are the last num_new
    new_sld_ids = sld_ids[-num_new:]
    old_sld_ids = sld_ids[:-num_new]

    # Insert after position 23 (0-indexed, which is slide 24 - the last content slide of Ch3)
    insert_pos = 24
    reordered = old_sld_ids[:insert_pos] + new_sld_ids + old_sld_ids[insert_pos:]

    # Clear and re-add
    for sld_id in list(sldIdLst):
        sldIdLst.remove(sld_id)

    for sld_id in reordered:
        sldIdLst.append(sld_id)

    print(f"Reordered: {len(reordered)} slides total, {num_new} new slides at position {insert_pos+1}-{insert_pos+num_new} (slide {insert_pos+2}-{insert_pos+num_new+1})")


# ========== MAIN ==========
print("=" * 60)
print("国有资产监管系统升级项目汇报 - 添加第三章亮点")
print("=" * 60)

# Copy original
print(f"\nCopying: {SRC_PATH} -> {DST_PATH}")
shutil.copy2(SRC_PATH, DST_PATH)

# Open
prs = Presentation(DST_PATH)
print(f"Original slides: {len(prs.slides)}")

# Create highlight slides
print(f"\nCreating {len(HIGHLIGHTS)} highlight slides...")
for i, hl in enumerate(HIGHLIGHTS):
    print(f"  [{i+1}/{len(HIGHLIGHTS)}] {hl['title'][:50]}...")
    create_highlight_slide(prs, hl)

# Reorder to place in Chapter 3
print(f"\nReordering slides to position highlights in Chapter 3...")
reorder_slides(prs, len(HIGHLIGHTS))

# Save
print(f"\nSaving presentation...")
prs.save(DST_PATH)

print(f"\n{'=' * 60}")
print(f"DONE! Output: {DST_PATH}")
print(f"Total slides: {len(prs.slides)} (original 33 + {len(HIGHLIGHTS)} new highlights)")
print(f"New slides inserted as slides 26-{25+len(HIGHLIGHTS)} in Chapter 3")
print(f"{'=' * 60}")
