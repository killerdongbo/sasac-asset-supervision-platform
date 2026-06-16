"""
Add 建设蓝图 (Construction Blueprint) architecture diagram slide into Chapter 2.
Inserted after slide 16 (建设目标), before slide 17 (Chapter 3 title).
Creates a layered architecture diagram framed as an overall blueprint.
"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8')

import shutil
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ========== PATHS ==========
SRC = r"d:\谢东波个人\智算\国资监管平台\国有资产监管系统升级项目汇报-V0.5.pptx"
DST = r"d:\谢东波个人\智算\国资监管平台\国有资产监管系统升级项目汇报-V0.6.pptx"

# ========== COLOR SCHEME (matching existing) ==========
BLUE_PRIMARY    = RGBColor(0x3C, 0x89, 0xFF)   # #3C89FF
BLUE_DARK       = RGBColor(0x1A, 0x56, 0xD6)   # #1A56D6
BLUE_DEEP       = RGBColor(0x0D, 0x3B, 0x9E)   # #0D3B9E  - deepest, for top layer
BLUE_MID        = RGBColor(0x25, 0x6B, 0xE6)   # #256BE6  - middle layer
BLUE_LIGHT      = RGBColor(0x5C, 0xA0, 0xFF)   # #5CA0FF  - light, for base layer
WHITE           = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK       = RGBColor(0x1A, 0x1A, 0x1A)
TEXT_GRAY       = RGBColor(0x66, 0x66, 0x66)
GREEN_CYAN      = RGBColor(0x00, 0xB8, 0xD4)   # cyan accent
ORANGE_ACCENT   = RGBColor(0xFF, 0x8F, 0x00)
PURPLE_ACCENT   = RGBColor(0x7C, 0x4D, 0xFF)
BORDER_GRAY     = RGBColor(0xE0, 0xE0, 0xE0)
BG_LIGHT_GRAY   = RGBColor(0xF5, 0xF8, 0xFF)
FONT_NAME       = '微软雅黑'

SLIDE_W = 12192000
SLIDE_H = 6858000

def add_textbox(slide, left, top, width, height, text, font_size=Pt(12),
                bold=False, color=TEXT_DARK, align=PP_ALIGN.LEFT, font_name=FONT_NAME,
                word_wrap=True):
    """Helper to add a text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.name = font_name
    run.font.color.rgb = color
    return txBox, tf, run

def add_rect(slide, left, top, width, height, fill_color, border_color=None, corner_radius=None):
    """Helper to add a rectangle with fill."""
    if corner_radius:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.adjustments[0] = corner_radius
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_arrow(slide, left, top, width, height, color=BLUE_PRIMARY):
    """Add a right arrow."""
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = color
    arrow.line.fill.background()
    return arrow

def create_blueprint_slide(prs):
    """Create the 建设蓝图 architecture diagram slide."""
    layout = prs.slide_layouts[4]  # '仅标题' - match Chapter 2 style
    slide = prs.slides.add_slide(layout)

    # White background
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE

    # ===== TOP TITLE AREA =====
    # Main title
    add_textbox(slide, Cm(8), Cm(0.4), Cm(18), Cm(0.8),
                '建设蓝图：分阶段达成，确保可见成效',
                Pt(26), bold=True, color=BLUE_DARK, align=PP_ALIGN.CENTER)

    # Subtitle
    add_textbox(slide, Cm(6), Cm(1.1), Cm(22), Cm(0.5),
                '以固定资产管理为起点，三阶段递进构建国资监管数字化体系',
                Pt(12), color=TEXT_GRAY, align=PP_ALIGN.CENTER)

    # Title underline
    add_rect(slide, Cm(12), Cm(0.95), Cm(10), Cm(0.03), BLUE_PRIMARY)

    # ===== OUTER FRAME (the "整体架构图框") =====
    frame_l = Cm(1.5)   # left
    frame_t = Cm(2.0)   # top
    frame_w = Cm(31.0)  # width
    frame_h = Cm(14.5)  # height

    # Outer frame dashed border
    frame = add_rect(slide, frame_l, frame_t, frame_w, frame_h,
                     RGBColor(0xFA, 0xFC, 0xFF), BORDER_GRAY, 0.03)
    frame.line.width = Pt(2)
    # Set dashed via XML
    ln = frame._element.find('.//' + qn('a:ln'))
    if ln is None:
        spPr = frame._element.find(qn('p:spPr'))
        if spPr is not None:
            ln = spPr.find(qn('a:ln'))
    if ln is not None:
        prstDash = etree.SubElement(ln, qn('a:prstDash'))
        prstDash.set('val', 'dash')

    # Frame label on top-left
    label_w = Cm(3.5)
    label = add_rect(slide, Cm(2.0), Cm(1.75), label_w, Cm(0.55), BLUE_PRIMARY, corner_radius=0.15)
    add_textbox(slide, Cm(2.0), Cm(1.75), label_w, Cm(0.55),
                '📋 整体建设蓝图', Pt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # ===== THREE LAYERED GOAL BLOCKS (nested, from base to top) =====
    # Design concept: three nested rounded rectangles, each containing the next

    # --- GOAL 1 (BASE / OUTERMOST) - 固定资产管理基础版本上线 ---
    g1_l = Cm(3.0);  g1_t = Cm(7.5);  g1_w = Cm(28.0);  g1_h = Cm(8.0)
    g1_bg = add_rect(slide, g1_l, g1_t, g1_w, g1_h,
                     RGBColor(0xE8, 0xF1, 0xFF), BLUE_LIGHT, 0.04)
    g1_bg.line.width = Pt(1.5)

    # G1 header bar
    g1_header = add_rect(slide, g1_l, g1_t, g1_w, Cm(1.2),
                         BLUE_LIGHT, corner_radius=0.04)
    # Clip top rounded corners by overlaying small rects? No, just accept slight visual diff
    # Actually rounded rect at top will have rounded top corners - that's fine for visual effect
    add_textbox(slide, Cm(3.5), g1_t + Cm(0.15), Cm(10), Cm(0.5),
                '目标一', Pt(11), bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_textbox(slide, Cm(3.5), g1_t + Cm(0.55), Cm(20), Cm(0.5),
                '固定资产管理基础版本上线 + 建立数据审核与上报流程',
                Pt(16), bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # G1 timeline
    add_rect(slide, Cm(23.0), g1_t + Cm(0.2), Cm(7.0), Cm(0.8),
             RGBColor(0xFF, 0xFF, 0xFF), corner_radius=0.2)
    add_textbox(slide, Cm(23.0), g1_t + Cm(0.25), Cm(7.0), Cm(0.7),
                '⏱ 2026年7月底', Pt(13), bold=True, color=BLUE_LIGHT, align=PP_ALIGN.CENTER)

    # G1 detail items (two columns)
    g1_items = [
        ('资产台账管理', '资产登记/编号/标签/变更记录'),
        ('资产流转管理', '入库/领用/调拨/报废全流程'),
        ('采购与折旧', '采购管理 + 财务折旧计提'),
        ('数据审核上报', '19类报表采集/审核/上报'),
        ('系统基础能力', '认证/权限RBAC/导入导出'),
        ('管理驾驶舱', '工作台/资产结构分析/动态'),
    ]
    col_w = Cm(8.5)
    for i, (item_title, item_desc) in enumerate(g1_items):
        row = i // 2
        col = i % 2
        x = g1_l + Cm(1.0) + col * (col_w + Cm(9.0))
        y = g1_t + Cm(1.6) + row * Cm(1.8)

        # Item icon
        icon = add_rect(slide, x, y, Cm(0.5), Cm(0.5), WHITE, BLUE_LIGHT, 1.0)
        add_textbox(slide, x, y, Cm(0.5), Cm(0.5),
                    '✓', Pt(9), bold=True, color=BLUE_LIGHT, align=PP_ALIGN.CENTER)
        # Item text
        add_textbox(slide, x + Cm(0.7), y - Cm(0.05), col_w, Cm(0.5),
                    item_title, Pt(11), bold=True, color=TEXT_DARK)
        add_textbox(slide, x + Cm(0.7), y + Cm(0.45), col_w, Cm(0.4),
                    item_desc, Pt(9), color=TEXT_GRAY)


    # --- GOAL 2 (MIDDLE) - 国企数字化平台 ---
    g2_l = Cm(5.0);  g2_t = Cm(3.8);  g2_w = Cm(24.0);  g2_h = Cm(3.2)
    g2_bg = add_rect(slide, g2_l, g2_t, g2_w, g2_h,
                     RGBColor(0xDB, 0xE7, 0xFF), BLUE_PRIMARY, 0.04)
    g2_bg.line.width = Pt(1.5)

    # G2 header bar
    g2_header = add_rect(slide, g2_l, g2_t, g2_w, Cm(1.15),
                         BLUE_PRIMARY, corner_radius=0.04)

    # Target label on header left
    add_textbox(slide, Cm(5.5), g2_t + Cm(0.1), Cm(8), Cm(0.5),
                '目标二', Pt(11), bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_textbox(slide, Cm(5.5), g2_t + Cm(0.5), Cm(22), Cm(0.5),
                '打造符合国资监管要求的国企数字化平台',
                Pt(16), bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # G2 timeline
    add_rect(slide, Cm(21.5), g2_t + Cm(0.15), Cm(6.5), Cm(0.8),
             WHITE, corner_radius=0.2)
    add_textbox(slide, Cm(21.5), g2_t + Cm(0.2), Cm(6.5), Cm(0.7),
                '⏱ 2026年底', Pt(13), bold=True, color=BLUE_PRIMARY, align=PP_ALIGN.CENTER)

    # G2 detail items (one row, 3 cols)
    g2_items = [
        ('全模块上线', '巡检/盘点/维保/审批\n工作流全面启用'),
        ('监管决策能力', '监管上报完整链路\n企业→集团→国资委'),
        ('平台化运营', '多租户SaaS模式\n全量国企统一接入'),
    ]
    for i, (item_title, item_desc) in enumerate(g2_items):
        x = g2_l + Cm(1.0) + i * Cm(7.8)
        y = g2_t + Cm(1.5)
        # Item card
        card = add_rect(slide, x, y, Cm(7.0), Cm(1.4), WHITE, BLUE_PRIMARY, 0.08)
        card.line.width = Pt(0.5)
        add_textbox(slide, x + Cm(0.4), y + Cm(0.15), Cm(6.2), Cm(0.5),
                    item_title, Pt(11), bold=True, color=BLUE_PRIMARY)
        add_textbox(slide, x + Cm(0.4), y + Cm(0.6), Cm(6.2), Cm(0.7),
                    item_desc, Pt(9), color=TEXT_GRAY)


    # --- GOAL 3 (TOP / INNERMOST) - 数据湖与大数据分析平台 ---
    g3_l = Cm(8.0);  g3_t = Cm(2.6);  g3_w = Cm(18.0);  g3_h = Cm(0.85)
    g3_header = add_rect(slide, g3_l, g3_t, g3_w, g3_h,
                         BLUE_DEEP, corner_radius=0.04)

    add_textbox(slide, Cm(8.5), g3_t + Cm(0.08), Cm(6), Cm(0.35),
                '目标三', Pt(10), bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_textbox(slide, Cm(8.5), g3_t + Cm(0.4), Cm(16), Cm(0.4),
                '建设国资数据湖与大数据分析平台',
                Pt(14), bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # G3 timeline
    add_rect(slide, Cm(20.0), g3_t + Cm(0.1), Cm(5.0), Cm(0.6),
             WHITE, corner_radius=0.2)
    add_textbox(slide, Cm(20.0), g3_t + Cm(0.12), Cm(5.0), Cm(0.55),
                '⏱ 2027年', Pt(11), bold=True, color=BLUE_DEEP, align=PP_ALIGN.CENTER)

    # ===== UP ARROW CONNECTORS between layers =====
    # Arrow from G1 to G2 (upward)
    arrow1 = slide.shapes.add_shape(MSO_SHAPE.UP_ARROW,
                                     Cm(16.0), g2_t + g2_h + Cm(0.1),
                                     Cm(2.0), Cm(0.55))
    arrow1.fill.solid()
    arrow1.fill.fore_color.rgb = BLUE_PRIMARY
    arrow1.line.fill.background()

    # Arrow from G2 to G3 (upward)
    arrow2 = slide.shapes.add_shape(MSO_SHAPE.UP_ARROW,
                                     Cm(16.0), g3_t + g3_h + Cm(0.05),
                                     Cm(2.0), Cm(0.55))
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = BLUE_PRIMARY
    arrow2.line.fill.background()

    # ===== SIDE ANNOTATIONS =====
    # Left side: phase labels
    phases = [
        (Cm(1.6), Cm(8.8), '夯实\n基础', BLUE_LIGHT),
        (Cm(1.6), Cm(5.2), '平台\n成型', BLUE_PRIMARY),
        (Cm(1.6), Cm(2.9), '数据\n智能', BLUE_DEEP),
    ]
    for px, py, ptext, pcolor in phases:
        badge = add_rect(slide, px, py, Cm(1.0), Cm(1.2), pcolor, corner_radius=0.3)
        add_textbox(slide, px, py + Cm(0.1), Cm(1.0), Cm(1.0),
                    ptext, Pt(10), bold=True, color=WHITE, align=PP_ALIGN.CENTER, font_name=FONT_NAME)

    # ===== BOTTOM TIMELINE =====
    timeline_y = Cm(15.8)
    # Timeline base line
    add_rect(slide, Cm(3.5), timeline_y + Cm(0.4), Cm(27.0), Cm(0.06), BLUE_PRIMARY)

    # Three milestone dots
    for i, (x_cm, label, desc) in enumerate([
        (Cm(4.5), '2026年7月', 'Phase 1\n基础上线'),
        (Cm(14.0), '2026年底', 'Phase 2\n平台成型'),
        (Cm(23.5), '2027年', 'Phase 3\n数据智能'),
    ]):
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x_cm, timeline_y + Cm(0.15), Cm(0.55), Cm(0.55))
        dot.fill.solid()
        colors = [BLUE_LIGHT, BLUE_PRIMARY, BLUE_DEEP]
        dot.fill.fore_color.rgb = colors[i]
        dot.line.fill.background()

        add_textbox(slide, x_cm - Cm(0.5), timeline_y + Cm(0.75), Cm(2.5), Cm(0.4),
                    label, Pt(10), bold=True, color=colors[i], align=PP_ALIGN.CENTER)
        add_textbox(slide, x_cm - Cm(0.5), timeline_y + Cm(1.05), Cm(2.5), Cm(0.5),
                    desc, Pt(8), color=TEXT_GRAY, align=PP_ALIGN.CENTER)

    # Phase connector arrows on timeline
    for x_start in [Cm(7.5), Cm(17.0)]:
        arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                      x_start, timeline_y + Cm(0.25),
                                      Cm(5.5), Cm(0.35))
        arr.fill.solid()
        arr.fill.fore_color.rgb = RGBColor(0xDD, 0xE8, 0xFF)
        arr.line.fill.background()

    # ===== FOOTER NOTE =====
    add_rect(slide, Cm(1.5), Cm(16.8), Cm(31.0), Cm(0.55),
             RGBColor(0xF0, 0xF5, 0xFF))
    add_textbox(slide, Cm(2.0), Cm(16.85), Cm(30.0), Cm(0.45),
                '建设策略：成熟软件匹配需求 → 以固定资产为核心做精做强 → 为企业预留可控空间 → 分阶段迭代逐步见效',
                Pt(9), color=TEXT_GRAY, align=PP_ALIGN.CENTER)

    print("  Blueprint slide created successfully.")
    return slide


def reorder_slides(prs, slide_to_move_idx, insert_after_idx):
    """
    Move a slide from slide_to_move_idx to right after insert_after_idx.
    Both are 0-indexed positions in the current sldIdLst order.
    """
    pres_part = prs.part
    pres_xml = pres_part._element
    sldIdLst = pres_xml.find(qn('p:sldIdLst'))
    sld_ids = list(sldIdLst.findall(qn('p:sldId')))

    # Remove the slide to move
    moved = sld_ids.pop(slide_to_move_idx)
    # Insert after the target position (adjust index since we removed one)
    if slide_to_move_idx < insert_after_idx:
        insert_after_idx -= 1
    sld_ids.insert(insert_after_idx + 1, moved)

    # Rebuild sldIdLst
    for sld_id in list(sldIdLst):
        sldIdLst.remove(sld_id)
    for sld_id in sld_ids:
        sldIdLst.append(sld_id)

    print(f"  Moved slide from position {slide_to_move_idx+1} to after position {insert_after_idx+1}")
    print(f"  Final slide count: {len(sld_ids)}")


# ========== MAIN ==========
print("=" * 60)
print("Adding 建设蓝图 slide to Chapter 2...")
print("=" * 60)

# Copy V0.5 -> V0.6
print(f"\nCopying: {SRC} -> {DST}")
shutil.copy2(SRC, DST)

prs = Presentation(DST)
print(f"Current slides: {len(prs.slides)}")

# Create the blueprint slide (added at end of presentation)
print(f"\nCreating blueprint slide...")
create_blueprint_slide(prs)

# Reorder: new slide is last (index 38, 0-indexed). Move to after slide 16 (index 15).
# Target position: after 建设目标 (slide 16, 0-indexed 15), before Chapter 3 title (slide 17)
new_slide_idx = len(prs.slides) - 1  # 38 (last)
insert_after = 15  # after slide 16 (0-indexed)
print(f"\nReordering: moving new slide (index {new_slide_idx}) to after slide {insert_after+1}...")
reorder_slides(prs, new_slide_idx, insert_after)

# Save
print(f"\nSaving presentation...")
prs.save(DST)

print(f"\n{'=' * 60}")
print(f"DONE! Output: {DST}")
print(f"Total slides: {len(prs.slides)}")
print(f"Blueprint slide inserted as slide 17 (after 建设目标, before Chapter 3)")
print(f"{'=' * 60}")
